"""Document text indexer — chunk + embed PDF/PPTx text for archive RAG chat.

For each doc_sources row:
  1. Extract per-page text (PyMuPDF for PDF, python-pptx slides for PPTX)
  2. Chunk ~500 tokens with ~50-token overlap, page number preserved
  3. Embed chunks with the BGE text model (passage mode, no query prefix)
  4. Persist doc_chunks rows, set index_status='ready' + chunk_count

Image extraction from the same uploads is handled separately by
routers/sources.py and is unaffected by this worker.
"""
from __future__ import annotations

import re
import uuid
from pathlib import Path
from typing import Any, Optional

import structlog

logger = structlog.get_logger()

# Token estimate: ~0.75 words per token for English prose → 500 tokens ≈ 375 words.
CHUNK_WORDS = 375
OVERLAP_WORDS = 38   # ≈ 50 tokens
MIN_CHUNK_CHARS = 40  # skip page fragments with no real content


def enqueue_doc_indexing(settings: Any, source_id: str) -> Optional[str]:
    """Best-effort RQ enqueue on the existing ingest queue. Returns job id or None."""
    if not settings.redis_url:
        return None
    try:
        import redis
        from rq import Queue

        q = Queue("ingest", connection=redis.from_url(settings.redis_url))
        job = q.enqueue(index_document, source_id, job_timeout=900)
        return job.id
    except Exception as exc:
        logger.warning("doc_index_enqueue_failed", source_id=source_id, error=str(exc))
        return None


def extract_pdf_pages(path: Path) -> list[tuple[int, str]]:
    """Return [(page_number, text)] for a PDF, 1-based pages."""
    import fitz  # PyMuPDF

    pages: list[tuple[int, str]] = []
    with fitz.open(str(path)) as doc:
        for i, page in enumerate(doc, start=1):
            pages.append((i, page.get_text("text") or ""))
    return pages


def extract_pptx_pages(path: Path) -> list[tuple[int, str]]:
    """Return [(slide_number, text)] for a PPTX, 1-based slides."""
    from pptx import Presentation

    pages: list[tuple[int, str]] = []
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(notes)
        pages.append((i, "\n".join(parts)))
    return pages


def chunk_pages(pages: list[tuple[int, str]]) -> list[tuple[int, int, str]]:
    """Split per-page text into overlapping word-window chunks.

    Returns [(page_number, chunk_index, text)] with chunk_index global per document.
    Chunks never span pages so citations stay exact.
    """
    chunks: list[tuple[int, int, str]] = []
    idx = 0
    for page_no, raw in pages:
        text = re.sub(r"[ \t]+", " ", raw or "").strip()
        if len(text) < MIN_CHUNK_CHARS:
            continue
        words = text.split()
        start = 0
        while start < len(words):
            window = words[start: start + CHUNK_WORDS]
            chunk_text = " ".join(window).strip()
            if len(chunk_text) >= MIN_CHUNK_CHARS:
                chunks.append((page_no, idx, chunk_text))
                idx += 1
            if start + CHUNK_WORDS >= len(words):
                break
            start += CHUNK_WORDS - OVERLAP_WORDS
    return chunks


def index_document(source_id: str) -> dict[str, Any]:
    """Extract, chunk, embed and persist text for one doc_sources row.

    Idempotent: existing chunks for the source are replaced.
    """
    import sqlalchemy as sa
    from sqlalchemy.orm import sessionmaker

    from app.config import get_settings
    from app.models.document import DocChunk, DocSource
    from app.services.text_embedder import embed_passages

    settings = get_settings()
    log = logger.bind(source_id=source_id)

    engine = sa.create_engine(settings.database_url, pool_pre_ping=True)
    Session = sessionmaker(bind=engine)

    with Session() as db:
        src = db.query(DocSource).filter(DocSource.id == uuid.UUID(source_id)).first()
        if src is None:
            log.error("doc_index_source_missing")
            return {"status": "error", "error": "source not found"}
        src.index_status = "indexing"
        src.index_error = None
        db.commit()
        file_type = src.file_type
        storage_path = Path(str(src.storage_path))

    try:
        if not storage_path.exists():
            raise FileNotFoundError(f"document file missing: {storage_path}")

        if file_type == "pdf":
            pages = extract_pdf_pages(storage_path)
        elif file_type == "pptx":
            pages = extract_pptx_pages(storage_path)
        else:
            raise ValueError(f"unsupported file_type: {file_type}")

        chunks = chunk_pages(pages)
        log.info("doc_index_chunked", pages=len(pages), chunks=len(chunks))

        # RQ forks per job, so this cap is scoped to this doc job only —
        # keeps BGE embedding from saturating CPU alongside image ingest jobs.
        try:
            import torch
            torch.set_num_threads(2)
        except Exception:
            pass

        embeddings = embed_passages([c[2] for c in chunks]) if chunks else None

        with Session() as db:
            db.query(DocChunk).filter(DocChunk.source_id == uuid.UUID(source_id)).delete()
            for i, (page_no, chunk_idx, text) in enumerate(chunks):
                db.add(DocChunk(
                    id=uuid.uuid4(),
                    source_id=uuid.UUID(source_id),
                    page_number=page_no,
                    chunk_index=chunk_idx,
                    text=text,
                    embedding=embeddings[i].tolist(),
                ))
            src = db.query(DocSource).filter(DocSource.id == uuid.UUID(source_id)).first()
            src.page_count = len(pages)
            src.chunk_count = len(chunks)
            src.index_status = "ready"
            src.index_error = None
            db.commit()

        log.info("doc_index_complete", chunks=len(chunks))
        return {"status": "ok", "source_id": source_id, "chunks": len(chunks)}

    except Exception as exc:
        import traceback
        err = f"{type(exc).__name__}: {exc}"
        log.error("doc_index_failed", error=err, traceback=traceback.format_exc())
        with Session() as db:
            src = db.query(DocSource).filter(DocSource.id == uuid.UUID(source_id)).first()
            if src is not None:
                src.index_status = "failed"
                src.index_error = err[:500]
                db.commit()
        return {"status": "error", "source_id": source_id, "error": err}
