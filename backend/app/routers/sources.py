"""Studio source ingestion endpoints.

Accepts:
  - URL  → scrape <img> tags, download, enqueue ingest
  - PDF  → extract embedded images via PyMuPDF, enqueue ingest
  - PPTX → extract embedded images via python-pptx, enqueue ingest
  - Video → sample up to 20 frames via opencv, enqueue ingest
  - S3   → list bucket objects, download images, enqueue ingest

Each ingest reuses `ingest_image()` from `workers.ingest_worker`, which
does CLIP embed + VLM artifact extraction.
"""
from __future__ import annotations

import asyncio
import hashlib
import io
import re
import tempfile
import uuid
from pathlib import Path
from typing import Optional
from urllib.parse import urljoin, urlparse

import httpx
import structlog
from fastapi import APIRouter, Depends, File, Form, HTTPException, Request, UploadFile
from pydantic import BaseModel, Field, HttpUrl
from sqlalchemy.orm import Session

from app.config import Settings, get_settings
from app.deps import get_db

logger = structlog.get_logger()

router = APIRouter(prefix="/studio/sources", tags=["studio-sources"])

# ── Limits ────────────────────────────────────────────────
MAX_IMG_BYTES = 12 * 1024 * 1024     # 12 MB per image
MAX_URL_IMAGES = 60                  # cap per URL page
MAX_PDF_IMAGES = 80                  # cap per PDF
MAX_PPTX_IMAGES = 80                 # cap per PPTX
MAX_VIDEO_FRAMES = 20                # spec'd by user
MAX_S3_IMAGES = 200                  # per call cap
ALLOWED_IMG_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tiff"}
ALLOWED_VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".webm"}
HTTP_TIMEOUT = httpx.Timeout(20.0, connect=8.0)
USER_AGENT = "Mozilla/5.0 (Visquery Studio Ingest) AppleWebKit/537.36"


# ── Schemas ───────────────────────────────────────────────
class URLIngestRequest(BaseModel):
    url: HttpUrl
    max_images: int = Field(default=MAX_URL_IMAGES, ge=1, le=MAX_URL_IMAGES)


class S3IngestRequest(BaseModel):
    bucket: str
    region: str = "us-east-1"
    access_key_id: str
    secret_access_key: str
    session_token: Optional[str] = None
    endpoint_url: Optional[str] = None       # for non-AWS S3 (B2, R2, MinIO)
    prefix: str = ""
    max_images: int = Field(default=MAX_S3_IMAGES, ge=1, le=MAX_S3_IMAGES)


class IngestResult(BaseModel):
    source_type: str
    discovered: int
    enqueued: int
    skipped: int
    job_ids: list[str] = []
    image_ids: list[str] = []
    errors: list[str] = []
    # Document text indexing (archive chat) — set for PDF/PPTx uploads only
    doc_source_id: Optional[str] = None
    doc_job_id: Optional[str] = None


# ── Helpers ───────────────────────────────────────────────
def _storage_root(settings: Settings) -> Path:
    root = Path(settings.storage_root) / "images"
    root.mkdir(parents=True, exist_ok=True)
    return root


def _persist_bytes(content: bytes, suffix: str, settings: Settings) -> tuple[Path, str]:
    """Write bytes to images dir with a fresh UUID name. Returns (path, sha256)."""
    if suffix.lower() not in ALLOWED_IMG_EXTS:
        suffix = ".jpg"
    sha = hashlib.sha256(content).hexdigest()
    name = f"{uuid.uuid4().hex}{suffix.lower()}"
    path = _storage_root(settings) / name
    path.write_bytes(content)
    return path, sha


def _is_image_bytes(content: bytes) -> Optional[str]:
    """Sniff magic bytes. Returns extension or None."""
    if len(content) < 12:
        return None
    if content[:3] == b"\xff\xd8\xff":
        return ".jpg"
    if content[:8] == b"\x89PNG\r\n\x1a\n":
        return ".png"
    if content[:6] in (b"GIF87a", b"GIF89a"):
        return ".gif"
    if content[:4] == b"RIFF" and content[8:12] == b"WEBP":
        return ".webp"
    if content[:2] == b"BM":
        return ".bmp"
    if content[:4] in (b"II*\x00", b"MM\x00*"):
        return ".tiff"
    return None


def _enqueue_ingest(
    settings: Settings,
    storage_path: str,
    source_url: str,
    source_title: str,
    spider_name: str,
    license_str: str = "user-uploaded",
) -> Optional[str]:
    """Enqueue ingest_image job via RQ. Returns job id or None."""
    try:
        import redis
        from rq import Queue
        from app.workers.ingest_worker import ingest_image

        q = Queue("ingest", connection=redis.from_url(settings.redis_url))
        job = q.enqueue(
            ingest_image,
            storage_path,
            source_url,
            source_title,
            license_str,
            spider_name,
            job_timeout=900,
        )
        return job.id
    except Exception as exc:
        logger.warning("studio_enqueue_failed", spider=spider_name, error=str(exc))
        return None


def _process_image_bytes(
    settings: Settings,
    content: bytes,
    source_url: str,
    source_title: str,
    spider_name: str,
    suffix: str = ".jpg",
) -> dict:
    """Persist + enqueue. Returns dict with status."""
    if len(content) > MAX_IMG_BYTES:
        return {"status": "skipped", "reason": "too_large"}
    ext = _is_image_bytes(content) or suffix
    if ext.lower() not in ALLOWED_IMG_EXTS:
        return {"status": "skipped", "reason": "not_image"}
    try:
        path, _ = _persist_bytes(content, ext, settings)
    except Exception as exc:
        return {"status": "error", "reason": str(exc)}
    job_id = _enqueue_ingest(
        settings,
        str(path.resolve()),
        source_url,
        source_title,
        spider_name,
    )
    return {"status": "queued" if job_id else "persisted_no_queue", "path": str(path), "job_id": job_id}


# ── URL scrape ────────────────────────────────────────────
@router.post("/url", response_model=IngestResult)
async def ingest_from_url(
    req: URLIngestRequest,
    settings: Settings = Depends(get_settings),
) -> IngestResult:
    """Scrape <img> tags from a webpage and ingest each image."""
    from bs4 import BeautifulSoup

    src_url = str(req.url)
    headers = {"User-Agent": USER_AGENT, "Accept": "text/html,*/*"}
    async with httpx.AsyncClient(timeout=HTTP_TIMEOUT, headers=headers, follow_redirects=True) as client:
        try:
            r = await client.get(src_url)
            r.raise_for_status()
        except Exception as exc:
            raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {exc}")

        soup = BeautifulSoup(r.text, "html.parser")
        page_title = (soup.title.string.strip() if soup.title and soup.title.string else src_url)

        # Collect candidate image URLs (img src, data-src, srcset largest)
        candidates: list[str] = []
        for img in soup.find_all("img"):
            for attr in ("src", "data-src", "data-original", "data-lazy-src"):
                v = img.get(attr)
                if v:
                    candidates.append(v.strip())
                    break
            sset = img.get("srcset")
            if sset:
                # last entry in srcset is usually largest
                parts = [p.strip().split(" ")[0] for p in sset.split(",") if p.strip()]
                if parts:
                    candidates.append(parts[-1])

        # Dedup + resolve absolute + filter ext
        seen = set()
        resolved: list[str] = []
        for c in candidates:
            absu = urljoin(src_url, c)
            ext = Path(urlparse(absu).path).suffix.lower()
            if ext and ext not in ALLOWED_IMG_EXTS:
                continue
            if absu in seen:
                continue
            seen.add(absu)
            resolved.append(absu)
            if len(resolved) >= req.max_images:
                break

        result = IngestResult(source_type="url", discovered=len(resolved), enqueued=0, skipped=0)

        # Download in parallel batches
        async def fetch_one(url: str) -> tuple[str, Optional[bytes], Optional[str]]:
            try:
                resp = await client.get(url)
                resp.raise_for_status()
                ct = resp.headers.get("content-type", "")
                if "image" not in ct and not any(url.lower().endswith(e) for e in ALLOWED_IMG_EXTS):
                    return url, None, "not_image_content_type"
                return url, resp.content, None
            except Exception as exc:
                return url, None, str(exc)

        # Batch concurrency
        sem = asyncio.Semaphore(6)
        async def bounded(u: str):
            async with sem:
                return await fetch_one(u)
        downloads = await asyncio.gather(*(bounded(u) for u in resolved))

        for url, content, err in downloads:
            if err or not content:
                result.skipped += 1
                if err:
                    result.errors.append(f"{url}: {err}")
                continue
            res = _process_image_bytes(
                settings, content, source_url=url, source_title=page_title, spider_name="studio_url",
                suffix=Path(urlparse(url).path).suffix.lower() or ".jpg",
            )
            if res["status"] == "queued":
                result.enqueued += 1
                if res.get("job_id"):
                    result.job_ids.append(res["job_id"])
            else:
                result.skipped += 1
                if res.get("reason"):
                    result.errors.append(f"{url}: {res['reason']}")

        return result


# ── Document text indexing (archive chat) ────────────────
def _register_document(
    content: bytes,
    filename: str,
    file_type: str,
    owner: Optional[str],
    settings: Settings,
    db: Session,
) -> tuple[Optional[str], Optional[str]]:
    """Persist the original file, create a doc_sources row, enqueue text indexing.

    Returns (doc_source_id, doc_job_id). Best-effort: failure here never blocks
    the image-extraction path.
    """
    try:
        from app.models.document import DocSource
        from app.workers.doc_indexer import enqueue_doc_indexing

        sha = hashlib.sha256(content).hexdigest()
        existing = db.query(DocSource).filter(DocSource.sha256 == sha).first()
        if existing:
            job_id = None
            if existing.index_status in ("queued", "failed"):
                job_id = enqueue_doc_indexing(settings, str(existing.id))
            return str(existing.id), job_id

        doc_dir = Path(settings.storage_root) / "documents"
        doc_dir.mkdir(parents=True, exist_ok=True)
        source_id = uuid.uuid4()
        ext = ".pdf" if file_type == "pdf" else ".pptx"
        dest = doc_dir / f"{source_id}{ext}"
        dest.write_bytes(content)

        db.add(DocSource(
            id=source_id,
            title=filename,
            file_type=file_type,
            storage_path=str(dest.resolve()),
            sha256=sha,
            index_status="queued",
            owner=owner,
        ))
        db.commit()

        job_id = enqueue_doc_indexing(settings, str(source_id))
        return str(source_id), job_id
    except Exception as exc:
        logger.warning("doc_register_failed", file=filename, error=str(exc))
        try:
            db.rollback()
        except Exception:
            pass
        return None, None


# ── PDF / PPTX upload ─────────────────────────────────────
def _extract_pdf_images(content: bytes) -> list[tuple[bytes, str]]:
    """Return list of (image_bytes, ext) from PDF via PyMuPDF (fitz)."""
    import fitz  # PyMuPDF
    out: list[tuple[bytes, str]] = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        for page in doc:
            for img_meta in page.get_images(full=True):
                xref = img_meta[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    if pix.alpha or pix.n > 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_bytes = pix.tobytes("png")
                    out.append((img_bytes, ".png"))
                    pix = None
                except Exception as exc:
                    logger.warning("pdf_extract_skip", xref=xref, error=str(exc))
                if len(out) >= MAX_PDF_IMAGES:
                    return out
    return out


def _extract_pptx_images(content: bytes) -> list[tuple[bytes, str]]:
    """Return list of (image_bytes, ext) from PPTX via python-pptx."""
    from pptx import Presentation
    from pptx.util import Emu  # noqa: F401
    out: list[tuple[bytes, str]] = []
    prs = Presentation(io.BytesIO(content))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13 and getattr(shape, "image", None):  # PICTURE
                try:
                    blob = shape.image.blob
                    ext = "." + shape.image.ext.lower()
                    out.append((blob, ext if ext in ALLOWED_IMG_EXTS else ".png"))
                except Exception as exc:
                    logger.warning("pptx_extract_skip", error=str(exc))
                if len(out) >= MAX_PPTX_IMAGES:
                    return out
    return out


@router.post("/pdf", response_model=IngestResult)
async def ingest_from_pdf(
    request: Request,
    file: UploadFile = File(...),
    settings: Settings = Depends(get_settings),
    db: Session = Depends(get_db),
) -> IngestResult:
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Empty file")
    if not (file.filename or "").lower().endswith(".pdf") and "pdf" not in (file.content_type or ""):
        raise HTTPException(status_code=400, detail="PDF file required")

    loop = asyncio.get_running_loop()
    try:
        images = await loop.run_in_executor(None, _extract_pdf_images, content)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"PDF parse failed: {exc}")

    result = IngestResult(source_type="pdf", discovered=len(images), enqueued=0, skipped=0)
    src_title = file.filename or "uploaded.pdf"

    # Archive text indexing — additive, image path below is unchanged
    owner = request.headers.get("X-Studio-Owner") or None
    result.doc_source_id, result.doc_job_id = _register_document(
        content, src_title, "pdf", owner, settings, db,
    )
    if result.doc_job_id:
        result.job_ids.append(result.doc_job_id)
    src_id = f"upload://pdf/{uuid.uuid4().hex}"
    for img_bytes, ext in images:
        res = _process_image_bytes(
            settings, img_bytes, source_url=src_id, source_title=src_title,
            spider_name="studio_pdf", suffix=ext,
        )
        if res["status"] == "queued":
            result.enqueued += 1
            if res.get("job_id"):
                result.job_ids.append(res["job_id"])
        else:
            result.skipped += 1
    return result


@router.post("/pptx", response_model=IngestResult)
async def ingest_from_pptx(
    request: Request,
    file: UploadFile = File(...),
    settings: Settings = Depends(get_settings),
    db: Session = Depends(get_db),
) -> IngestResult:
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Empty file")
    name = (file.filename or "").lower()
    if not (name.endswith(".pptx") or name.endswith(".ppt")):
        raise HTTPException(status_code=400, detail="PPTX file required")

    loop = asyncio.get_running_loop()
    try:
        images = await loop.run_in_executor(None, _extract_pptx_images, content)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"PPTX parse failed: {exc}")

    result = IngestResult(source_type="pptx", discovered=len(images), enqueued=0, skipped=0)
    src_title = file.filename or "uploaded.pptx"

    # Archive text indexing — additive, image path below is unchanged
    owner = request.headers.get("X-Studio-Owner") or None
    result.doc_source_id, result.doc_job_id = _register_document(
        content, src_title, "pptx", owner, settings, db,
    )
    if result.doc_job_id:
        result.job_ids.append(result.doc_job_id)
    src_id = f"upload://pptx/{uuid.uuid4().hex}"
    for img_bytes, ext in images:
        res = _process_image_bytes(
            settings, img_bytes, source_url=src_id, source_title=src_title,
            spider_name="studio_pptx", suffix=ext,
        )
        if res["status"] == "queued":
            result.enqueued += 1
            if res.get("job_id"):
                result.job_ids.append(res["job_id"])
        else:
            result.skipped += 1
    return result


# ── Video frame extract ───────────────────────────────────
def _extract_video_frames(content: bytes, suffix: str, max_frames: int = MAX_VIDEO_FRAMES) -> list[bytes]:
    """Sample evenly-spaced frames from a video. Returns list of JPEG bytes."""
    import cv2
    import numpy as np

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        cap = cv2.VideoCapture(tmp_path)
        if not cap.isOpened():
            raise RuntimeError("Cannot open video")
        total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT)) or 0
        if total <= 0:
            # Stream: read frames sequentially up to cap
            frames: list[bytes] = []
            idx = 0
            while len(frames) < max_frames:
                ok, frame = cap.read()
                if not ok:
                    break
                if idx % 30 == 0:  # ~1 fps assuming 30fps
                    ok2, buf = cv2.imencode(".jpg", frame, [int(cv2.IMWRITE_JPEG_QUALITY), 88])
                    if ok2:
                        frames.append(buf.tobytes())
                idx += 1
            cap.release()
            return frames

        # Even sampling
        step = max(1, total // max_frames)
        frames: list[bytes] = []
        for i in range(max_frames):
            cap.set(cv2.CAP_PROP_POS_FRAMES, min(i * step, total - 1))
            ok, frame = cap.read()
            if not ok or frame is None:
                continue
            ok2, buf = cv2.imencode(".jpg", frame, [int(cv2.IMWRITE_JPEG_QUALITY), 88])
            if ok2:
                frames.append(buf.tobytes())
        cap.release()
        return frames
    finally:
        try:
            Path(tmp_path).unlink(missing_ok=True)
        except Exception:
            pass


@router.post("/video", response_model=IngestResult)
async def ingest_from_video(
    file: UploadFile = File(...),
    max_frames: int = Form(MAX_VIDEO_FRAMES),
    settings: Settings = Depends(get_settings),
) -> IngestResult:
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Empty file")
    suffix = Path(file.filename or "video.mp4").suffix.lower()
    if suffix not in ALLOWED_VIDEO_EXTS:
        raise HTTPException(status_code=400, detail=f"Video extension not supported: {suffix}")

    max_frames = max(1, min(int(max_frames), MAX_VIDEO_FRAMES))

    loop = asyncio.get_running_loop()
    try:
        frames = await loop.run_in_executor(None, _extract_video_frames, content, suffix, max_frames)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Video parse failed: {exc}")

    result = IngestResult(source_type="video", discovered=len(frames), enqueued=0, skipped=0)
    src_title = file.filename or "uploaded.mp4"
    src_id = f"upload://video/{uuid.uuid4().hex}"
    for frame_bytes in frames:
        res = _process_image_bytes(
            settings, frame_bytes, source_url=src_id, source_title=src_title,
            spider_name="studio_video", suffix=".jpg",
        )
        if res["status"] == "queued":
            result.enqueued += 1
            if res.get("job_id"):
                result.job_ids.append(res["job_id"])
        else:
            result.skipped += 1
    return result


# ── S3 bucket ingest ──────────────────────────────────────
@router.post("/s3", response_model=IngestResult)
async def ingest_from_s3(
    req: S3IngestRequest,
    settings: Settings = Depends(get_settings),
) -> IngestResult:
    """List a bucket prefix, download images, enqueue ingest.

    NOTE: credentials are used only for this call and never persisted.
    """
    import boto3
    from botocore.config import Config
    from botocore.exceptions import BotoCoreError, ClientError

    cfg = Config(retries={"max_attempts": 3, "mode": "standard"})
    client_kwargs = {
        "aws_access_key_id": req.access_key_id,
        "aws_secret_access_key": req.secret_access_key,
        "region_name": req.region,
        "config": cfg,
    }
    if req.session_token:
        client_kwargs["aws_session_token"] = req.session_token
    if req.endpoint_url:
        client_kwargs["endpoint_url"] = req.endpoint_url

    loop = asyncio.get_running_loop()
    try:
        s3 = boto3.client("s3", **client_kwargs)

        def _list_keys() -> list[dict]:
            keys: list[dict] = []
            paginator = s3.get_paginator("list_objects_v2")
            for page in paginator.paginate(Bucket=req.bucket, Prefix=req.prefix or ""):
                for obj in page.get("Contents", []) or []:
                    k = obj["Key"]
                    if Path(k).suffix.lower() in ALLOWED_IMG_EXTS:
                        keys.append({"Key": k, "Size": obj.get("Size", 0)})
                    if len(keys) >= req.max_images:
                        return keys
            return keys

        keys = await loop.run_in_executor(None, _list_keys)
    except (BotoCoreError, ClientError) as exc:
        raise HTTPException(status_code=400, detail=f"S3 list failed: {exc}")
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"S3 client error: {exc}")

    result = IngestResult(source_type="s3", discovered=len(keys), enqueued=0, skipped=0)
    src_title = f"s3://{req.bucket}/{req.prefix}".rstrip("/")

    def _download(key: str) -> bytes:
        buf = io.BytesIO()
        s3.download_fileobj(req.bucket, key, buf)
        return buf.getvalue()

    for obj in keys:
        key = obj["Key"]
        if obj.get("Size", 0) > MAX_IMG_BYTES:
            result.skipped += 1
            result.errors.append(f"{key}: too_large")
            continue
        try:
            content = await loop.run_in_executor(None, _download, key)
        except Exception as exc:
            result.skipped += 1
            result.errors.append(f"{key}: {exc}")
            continue
        res = _process_image_bytes(
            settings, content,
            source_url=f"s3://{req.bucket}/{key}",
            source_title=src_title,
            spider_name="studio_s3",
            suffix=Path(key).suffix.lower() or ".jpg",
        )
        if res["status"] == "queued":
            result.enqueued += 1
            if res.get("job_id"):
                result.job_ids.append(res["job_id"])
        else:
            result.skipped += 1

    return result


# ── Job status batch ──────────────────────────────────────
class JobStatusRequest(BaseModel):
    job_ids: list[str]


@router.post("/jobs/status")
async def jobs_status(
    req: JobStatusRequest,
    settings: Settings = Depends(get_settings),
) -> dict:
    import redis
    from rq.job import Job
    from rq.exceptions import NoSuchJobError

    r = redis.from_url(settings.redis_url)
    out: dict[str, dict] = {}
    for jid in req.job_ids[:200]:
        try:
            j = Job.fetch(jid, connection=r)
            out[jid] = {
                "status": j.get_status().value,
                "result": j.result if j.is_finished else None,
            }
        except NoSuchJobError:
            out[jid] = {"status": "missing"}
        except Exception as exc:
            out[jid] = {"status": "error", "error": str(exc)}
    return {"jobs": out}
